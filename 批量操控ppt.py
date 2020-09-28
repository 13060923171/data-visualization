from pptx import Presentation
from datetime import datetime
from openpyxl import load_workbook

prs = Presentation('奖学金证书模板.pptx')
title_slide_layout = prs.slide_layouts[0]
slide = prs.slide.add_slide(title_slide_layout)
winner_name = slide.placeholders[0]
certificate_type = slide.placeholders[1]
this_is_to_certify_that = slide.placeholders[20]
winning_reason = slide.placeholders[19]
award_presenter = slide.placeholders[17]
award_data = slide.placeholders[21]
wb = load_workbook(filename='学生名单及获奖理由.xlsx')
sheet = wb.active
for row in sheet.iter_rows(min_row=2):
    name,reason = row
    winner_name.text = name.value
    certificate_type.text = reason.value
    this_is_to_certify_that.text = '兹证明'
    winning_reason.text = '思想品德好、学习好、身体健康'
    award_presenter.text = '曾老师'
    award_data.text = str(datetime.now().date())
    prs.save(f'./output/{name.value}.pptx')